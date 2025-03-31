import os
import re
import logging
from tempfile import NamedTemporaryFile
from typing import Optional, Dict, List
from pathlib import Path
from PyPDF2 import PdfReader
from docx import Document
import pptx
import openpyxl  # For Excel files
import pytesseract
from PIL import Image, UnidentifiedImageError
import pandas as pd  # For CSV and Excel
import pdfplumber  # Alternative PDF parser
from app.utils.config import settings
from fastapi import HTTPException, UploadFile

logger = logging.getLogger(__name__)

class FileParser:
    def __init__(self):
        """Initialize the file parser with OCR configuration"""
        self._configure_ocr()
        self.supported_formats = {
            '.pdf': self._parse_pdf,
            '.docx': self._parse_docx,
            '.doc': self._parse_doc,
            '.pptx': self._parse_pptx,
            '.xlsx': self._parse_excel,
            '.csv': self._parse_csv,
            '.txt': self._parse_text,
            '.png': self._parse_image,
            '.jpg': self._parse_image,
            '.jpeg': self._parse_image,
        }
        self.max_file_size = settings.MAX_FILE_SIZE_MB * 1024 * 1024  # Convert MB to bytes

    def _configure_ocr(self):
        """Configure Tesseract OCR path if specified"""
        if hasattr(settings, 'TESSERACT_PATH') and settings.TESSERACT_PATH:
            pytesseract.pytesseract.tesseract_cmd = settings.TESSERACT_PATH
        elif os.name == 'nt':  # Windows default path
            default_path = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
            if os.path.exists(default_path):
                pytesseract.pytesseract.tesseract_cmd = default_path

    async def parse_file(self, file: UploadFile) -> Dict[str, str]:
        """
        Parse content from various file formats with validation
        
        Args:
            file: FastAPI UploadFile object
            
        Returns:
            Dictionary containing:
            - 'content': Extracted text
            - 'metadata': File metadata
            
        Raises:
            HTTPException: For unsupported formats or parsing errors
        """
        try:
            # Validate file
            self._validate_file(file)
            
            # Get file extension
            file_ext = Path(file.filename).suffix.lower()
            
            if file_ext not in self.supported_formats:
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported file format: {file_ext}"
                )
            
            # Use temp file for parsing
            with NamedTemporaryFile(delete=False) as temp_file:
                content = await file.read()
                temp_file.write(content)
                temp_path = temp_file.name
            
            try:
                # Parse the file
                parse_func = self.supported_formats[file_ext]
                result = parse_func(temp_path)
                
                # Clean up temp file
                os.unlink(temp_path)
                
                return {
                    'content': self._clean_text(result['content']),
                    'metadata': {
                        **result['metadata'],
                        'file_name': file.filename,
                        'file_size': len(content),
                        'file_type': file_ext[1:].upper()  # Remove dot
                    }
                }
            except Exception as e:
                os.unlink(temp_path)
                raise e
                
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"Error parsing file {file.filename}: {str(e)}")
            raise HTTPException(
                status_code=500,
                detail=f"Failed to parse file: {str(e)}"
            )

    def _validate_file(self, file: UploadFile):
        """Validate file before processing"""
        if not file.filename:
            raise HTTPException(
                status_code=400,
                detail="No filename provided"
            )
            
        if file.size > self.max_file_size:
            raise HTTPException(
                status_code=413,
                detail=f"File too large. Max size is {settings.MAX_FILE_SIZE_MB}MB"
            )

    def _parse_pdf(self, file_path: str) -> Dict[str, str]:
        """Parse text from PDF file using multiple methods"""
        content = ""
        metadata = {}
        
        try:
            # First try with pdfplumber for better text extraction
            with pdfplumber.open(file_path) as pdf:
                content = "\n".join(
                    page.extract_text() or "" 
                    for page in pdf.pages
                )
                metadata['page_count'] = len(pdf.pages)
                
            # Fallback to PyPDF2 if no text extracted
            if not content.strip():
                with open(file_path, 'rb') as f:
                    reader = PdfReader(f)
                    content = "\n".join(
                        page.extract_text() or "" 
                        for page in reader.pages
                    )
                    metadata['page_count'] = len(reader.pages)
                    
        except Exception as e:
            logger.warning(f"PDF parsing error: {str(e)}")
            raise HTTPException(
                status_code=400,
                detail="Failed to parse PDF file"
            )
            
        return {
            'content': content,
            'metadata': metadata
        }

    def _parse_docx(self, file_path: str) -> Dict[str, str]:
        """Parse text from DOCX file"""
        try:
            doc = Document(file_path)
            paragraphs = [para.text for para in doc.paragraphs]
            
            # Extract tables
            tables = []
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        tables.append(cell.text)
            
            content = "\n".join(paragraphs + tables)
            
            return {
                'content': content,
                'metadata': {
                    'paragraph_count': len(paragraphs),
                    'table_count': len(doc.tables)
                }
            }
        except Exception as e:
            logger.warning(f"DOCX parsing error: {str(e)}")
            raise HTTPException(
                status_code=400,
                detail="Failed to parse DOCX file"
            )

    def _parse_doc(self, file_path: str) -> Dict[str, str]:
        """Parse text from legacy DOC format (requires antiword)"""
        try:
            # Try using antiword if available
            if os.system("which antiword > /dev/null") == 0:
                import subprocess
                result = subprocess.run(
                    ["antiword", file_path],
                    capture_output=True,
                    text=True
                )
                if result.returncode == 0:
                    return {
                        'content': result.stdout,
                        'metadata': {'format': 'DOC'}
                    }
            
            # Fallback to textract if installed
            try:
                import textract
                content = textract.process(file_path).decode('utf-8')
                return {
                    'content': content,
                    'metadata': {'format': 'DOC'}
                }
            except ImportError:
                pass
                
            raise HTTPException(
                status_code=400,
                detail="DOC parsing requires antiword or textract installation"
            )
        except Exception as e:
            logger.warning(f"DOC parsing error: {str(e)}")
            raise HTTPException(
                status_code=400,
                detail="Failed to parse DOC file"
            )

    def _parse_pptx(self, file_path: str) -> Dict[str, str]:
        """Parse text from PowerPoint files"""
        try:
            presentation = pptx.Presentation(file_path)
            content = []
            
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content.append(shape.text)
            
            return {
                'content': "\n".join(content),
                'metadata': {
                    'slide_count': len(presentation.slides)
                }
            }
        except Exception as e:
            logger.warning(f"PPTX parsing error: {str(e)}")
            raise HTTPException(
                status_code=400,
                detail="Failed to parse PowerPoint file"
            )

    def _parse_excel(self, file_path: str) -> Dict[str, str]:
        """Parse text from Excel files"""
        try:
            # Read all sheets
            excel_data = pd.read_excel(file_path, sheet_name=None)
            content = []
            
            for sheet_name, df in excel_data.items():
                content.append(f"=== Sheet: {sheet_name} ===")
                content.append(df.to_string())
            
            return {
                'content': "\n".join(content),
                'metadata': {
                    'sheet_count': len(excel_data),
                    'columns': list(excel_data.values())[0].columns.tolist() if excel_data else []
                }
            }
        except Exception as e:
            logger.warning(f"Excel parsing error: {str(e)}")
            raise HTTPException(
                status_code=400,
                detail="Failed to parse Excel file"
            )

    def _parse_csv(self, file_path: str) -> Dict[str, str]:
        """Parse text from CSV files"""
        try:
            df = pd.read_csv(file_path)
            return {
                'content': df.to_string(),
                'metadata': {
                    'row_count': len(df),
                    'columns': df.columns.tolist()
                }
            }
        except Exception as e:
            logger.warning(f"CSV parsing error: {str(e)}")
            raise HTTPException(
                status_code=400,
                detail="Failed to parse CSV file"
            )

    def _parse_text(self, file_path: str) -> Dict[str, str]:
        """Parse plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                
            return {
                'content': content,
                'metadata': {
                    'line_count': len(content.splitlines())
                }
            }
        except Exception as e:
            logger.warning(f"Text file parsing error: {str(e)}")
            raise HTTPException(
                status_code=400,
                detail="Failed to parse text file"
            )

    def _parse_image(self, file_path: str) -> Dict[str, str]:
        """Parse text from images using OCR"""
        try:
            image = Image.open(file_path)
            
            # Preprocess image for better OCR results
            image = image.convert('L')  # Convert to grayscale
            content = pytesseract.image_to_string(image)
            
            return {
                'content': content,
                'metadata': {
                    'dimensions': f"{image.width}x{image.height}",
                    'mode': image.mode
                }
            }
        except UnidentifiedImageError:
            raise HTTPException(
                status_code=400,
                detail="Invalid image file"
            )
        except Exception as e:
            logger.warning(f"Image OCR error: {str(e)}")
            raise HTTPException(
                status_code=400,
                detail="Failed to extract text from image"
            )

    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        if not text:
            return ""
            
        # Remove excessive whitespace and non-printable characters
        text = re.sub(r'\s+', ' ', text).strip()
        text = ''.join(char for char in text if char.isprintable())
        
        # Normalize line endings and remove special characters
        text = text.replace('\r\n', '\n').replace('\r', '\n')
        text = re.sub(r'[\x00-\x1f\x7f-\x9f]', '', text)
        
        return text

    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        return list(self.supported_formats.keys())

    async def extract_metadata(self, file: UploadFile) -> Dict:
        """Extract basic metadata from file without full parsing"""
        try:
            file_ext = Path(file.filename).suffix.lower()
            
            if file_ext == '.pdf':
                with open(file.filename, 'rb') as f:
                    reader = PdfReader(f)
                    return {
                        'page_count': len(reader.pages),
                        'author': reader.metadata.get('/Author', ''),
                        'title': reader.metadata.get('/Title', '')
                    }
                    
            elif file_ext in ('.docx', '.doc'):
                # Similar metadata extraction for other formats
                pass
                
            return {}
            
        except Exception as e:
            logger.error(f"Metadata extraction failed: {str(e)}")
            return {}
